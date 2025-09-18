#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AICoreDirector 增强版PPT生成器
生成更加丰富、美观的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_enhanced_ppt():
    """创建增强版AICoreDirector演示文稿"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 定义主题色彩
    primary_color = RGBColor(41, 128, 185)      # 主色调：蓝色
    secondary_color = RGBColor(52, 152, 219)    # 次要色：浅蓝
    accent_color = RGBColor(231, 76, 60)        # 强调色：红色
    success_color = RGBColor(46, 204, 113)      # 成功色：绿色
    warning_color = RGBColor(241, 196, 15)      # 警告色：黄色
    text_color = RGBColor(44, 62, 80)           # 文字色：深灰
    
    # 1. 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加背景渐变效果
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.gradient()
    background.fill.gradient.gradient_stops[0].color.rgb = RGBColor(236, 240, 241)
    background.fill.gradient.gradient_stops[1].color.rgb = RGBColor(189, 195, 199)
    background.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "AICoreDirector"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "企业级AI能力管理中枢平台"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = secondary_color
    
    # 版本信息
    version_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(0.8))
    version_frame = version_box.text_frame
    version_frame.text = "统一管理 | 智能调度 | 灵活扩展 | 成本可控"
    version_para = version_frame.paragraphs[0]
    version_para.alignment = PP_ALIGN.CENTER
    version_para.font.size = Pt(20)
    version_para.font.color.rgb = text_color
    
    # 2. 目录页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 249, 250)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 目录内容
    toc_items = [
        "1. 项目背景与痛点分析",
        "2. 解决方案与核心架构",
        "3. 平台四大支柱能力",
        "4. 技术架构与实现",
        "5. 核心功能演示",
        "6. 商业价值与竞争优势",
        "7. 发展规划与路线图"
    ]
    
    for i, item in enumerate(toc_items):
        item_box = slide.shapes.add_textbox(Inches(1.5), Inches(2 + i * 0.6), Inches(10.83), Inches(0.5))
        item_frame = item_box.text_frame
        item_frame.text = item
        item_para = item_frame.paragraphs[0]
        item_para.font.size = Pt(24)
        item_para.font.color.rgb = text_color
        item_para.font.bold = True
    
    # 3. 背景与痛点分析
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "背景与痛点分析"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 痛点分析图表
    pain_points = [
        ("模型管理复杂", "多模型、多版本、多环境难以统一管理"),
        ("资源调度低效", "缺乏智能路由，无法根据业务需求自动选择最优模型"),
        ("成本控制困难", "缺乏统一的成本监控和优化机制"),
        ("运维监控缺失", "模型服务状态、性能指标缺乏实时监控"),
        ("扩展性不足", "新模型接入困难，缺乏标准化接口")
    ]
    
    # 创建痛点分析图表
    for i, (title, desc) in enumerate(pain_points):
        # 痛点标题
        pain_title = slide.shapes.add_textbox(
            Inches(0.8 + i * 2.4), Inches(1.5), Inches(2), Inches(0.6)
        )
        pain_title_frame = pain_title.text_frame
        pain_title_frame.text = title
        pain_title_para = pain_title_frame.paragraphs[0]
        pain_title_para.font.size = Pt(18)
        pain_title_para.font.bold = True
        pain_title_para.font.color.rgb = accent_color
        pain_title_para.alignment = PP_ALIGN.CENTER
        
        # 痛点描述
        pain_desc = slide.shapes.add_textbox(
            Inches(0.8 + i * 2.4), Inches(2.2), Inches(2), Inches(1.2)
        )
        pain_desc_frame = pain_desc.text_frame
        pain_desc_frame.text = desc
        pain_desc_para = pain_desc_frame.paragraphs[0]
        pain_desc_para.font.size = Pt(14)
        pain_desc_para.font.color.rgb = text_color
        pain_desc_para.alignment = PP_ALIGN.CENTER
        pain_desc_para.word_wrap = True
        
        # 添加图标背景
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 2.4), Inches(1.5), Inches(2), Inches(2)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        icon_bg.line.color.rgb = RGBColor(189, 195, 199)
        icon_bg.line.width = Pt(2)
        icon_bg.shadow.inherit = False
    
    # 4. 解决方案与核心架构
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "解决方案与核心架构"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 解决方案描述
    solution_text = """
    AICoreDirector作为企业级AI能力管理中枢，提供统一的模型服务管理平台：
    
    • 统一接入：支持LLM、传统ML模型、自研AI组件的统一接入
    • 智能调度：基于业务标签、性能指标、成本偏好的智能路由
    • 灵活扩展：插件化架构，支持热插拔和动态扩展
    • 成本可控：实时监控、成本分析、资源优化
    """
    
    solution_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(4))
    solution_frame = solution_box.text_frame
    solution_frame.text = solution_text
    solution_para = solution_frame.paragraphs[0]
    solution_para.font.size = Pt(18)
    solution_para.font.color.rgb = text_color
    solution_para.font.bold = False
    
    # 架构图占位符
    arch_placeholder = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(4))
    arch_frame = arch_placeholder.text_frame
    arch_frame.text = "[架构图占位符]"
    arch_para = arch_frame.paragraphs[0]
    arch_para.font.size = Pt(24)
    arch_para.font.color.rgb = RGBColor(128, 128, 128)
    arch_para.alignment = PP_ALIGN.CENTER
    
    # 5. 平台四大支柱能力
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "平台四大支柱能力"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 四大支柱
    pillars = [
        ("智能路由调度", "基于多维度指标的智能模型选择", success_color),
        ("插件化架构", "热插拔、动态扩展的业务能力", primary_color),
        ("统一监控治理", "全链路监控、成本分析、性能优化", warning_color),
        ("标准化接口", "RESTful API、SDK、多语言支持", accent_color)
    ]
    
    # 创建支柱图表
    for i, (title, desc, color) in enumerate(pillars):
        # 支柱背景
        pillar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 2.8), Inches(1.8), Inches(2.4), Inches(3.5)
        )
        pillar_bg.fill.solid()
        pillar_bg.fill.fore_color.rgb = color
        pillar_bg.line.color.rgb = color
        pillar_bg.line.width = Pt(3)
        pillar_bg.shadow.inherit = False
        
        # 支柱标题
        pillar_title = slide.shapes.add_textbox(
            Inches(0.9 + i * 2.8), Inches(2), Inches(2.2), Inches(0.8)
        )
        pillar_title_frame = pillar_title.text_frame
        pillar_title_frame.text = title
        pillar_title_para = pillar_title_frame.paragraphs[0]
        pillar_title_para.font.size = Pt(20)
        pillar_title_para.font.bold = True
        pillar_title_para.font.color.rgb = RGBColor(255, 255, 255)
        pillar_title_para.alignment = PP_ALIGN.CENTER
        
        # 支柱描述
        pillar_desc = slide.shapes.add_textbox(
            Inches(0.9 + i * 2.8), Inches(2.9), Inches(2.2), Inches(2.2)
        )
        pillar_desc_frame = pillar_desc.text_frame
        pillar_desc_frame.text = desc
        pillar_desc_para = pillar_desc_frame.paragraphs[0]
        pillar_desc_para.font.size = Pt(16)
        pillar_desc_para.font.color.rgb = RGBColor(255, 255, 255)
        pillar_desc_para.alignment = PP_ALIGN.CENTER
        pillar_desc_para.word_wrap = True
    
    # 6. 技术架构与实现
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "技术架构与实现"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 技术栈
    tech_stack = [
        ("后端框架", "FastAPI + Python 3.8+", "高性能异步Web框架"),
        ("前端技术", "Vue3 + Element Plus", "现代化响应式UI"),
        ("数据存储", "JSON + 文件系统", "轻量级配置管理"),
        ("部署方式", "Docker + K8s", "容器化微服务架构"),
        ("监控告警", "Prometheus + Grafana", "全方位监控体系")
    ]
    
    for i, (category, tech, desc) in enumerate(tech_stack):
        # 分类标题
        cat_title = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5 + i * 0.8), Inches(2.5), Inches(0.6)
        )
        cat_title_frame = cat_title.text_frame
        cat_title_frame.text = category
        cat_title_para = cat_title_frame.paragraphs[0]
        cat_title_para.font.size = Pt(18)
        cat_title_para.font.bold = True
        cat_title_para.font.color.rgb = primary_color
        
        # 技术名称
        tech_name = slide.shapes.add_textbox(
            Inches(3.5), Inches(1.5 + i * 0.8), Inches(4), Inches(0.6)
        )
        tech_name_frame = tech_name.text_frame
        tech_name_frame.text = tech
        tech_name_para = tech_name_frame.paragraphs[0]
        tech_name_para.font.size = Pt(18)
        tech_name_para.font.color.rgb = text_color
        
        # 技术描述
        tech_desc = slide.shapes.add_textbox(
            Inches(8), Inches(1.5 + i * 0.8), Inches(4.5), Inches(0.6)
        )
        tech_desc_frame = tech_desc.text_frame
        tech_desc_frame.text = desc
        tech_desc_para = tech_desc_frame.paragraphs[0]
        tech_desc_para.font.size = Pt(16)
        tech_desc_para.font.color.rgb = RGBColor(128, 128, 128)
    
    # 7. 核心功能演示
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "核心功能演示"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 功能特性
    features = [
        ("多模型管理", "支持GPT、Claude、GLM、文心一言等多种LLM模型"),
        ("智能路由", "基于标签、性能、成本的智能模型选择"),
        ("插件系统", "热插拔业务插件，支持批量并发调用"),
        ("服务注册", "外部AI/ML服务动态注册与发现"),
        ("Prompt管理", "集中化提示词模板管理，支持多种格式"),
        ("监控治理", "实时QPS、成本、健康状态监控")
    ]
    
    # 创建功能网格
    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        
        # 功能卡片背景
        feature_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + col * 6.5), Inches(1.5 + row * 2.2), Inches(5.8), Inches(1.8)
        )
        feature_bg.fill.solid()
        feature_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        feature_bg.line.color.rgb = RGBColor(189, 195, 199)
        feature_bg.line.width = Pt(2)
        feature_bg.shadow.inherit = False
        
        # 功能标题
        feature_title = slide.shapes.add_textbox(
            Inches(1 + col * 6.5), Inches(1.7 + row * 2.2), Inches(5.4), Inches(0.6)
        )
        feature_title_frame = feature_title.text_frame
        feature_title_frame.text = title
        feature_title_para = feature_title_frame.paragraphs[0]
        feature_title_para.font.size = Pt(20)
        feature_title_para.font.bold = True
        feature_title_para.font.color.rgb = primary_color
        
        # 功能描述
        feature_desc = slide.shapes.add_textbox(
            Inches(1 + col * 6.5), Inches(2.4 + row * 2.2), Inches(5.4), Inches(0.8)
        )
        feature_desc_frame = feature_desc.text_frame
        feature_desc_frame.text = desc
        feature_desc_para = feature_desc_frame.paragraphs[0]
        feature_desc_para.font.size = Pt(16)
        feature_desc_para.font.color.rgb = text_color
        feature_desc_para.word_wrap = True
    
    # 8. 商业价值与竞争优势
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "商业价值与竞争优势"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 商业价值
    business_value = [
        ("成本优化", "统一管理降低30%运维成本", "智能调度提升资源利用率"),
        ("效率提升", "标准化接口提升开发效率", "插件化架构加速业务迭代"),
        ("风险控制", "统一监控降低服务风险", "标准化治理提升合规性"),
        ("扩展能力", "支持多种AI模型接入", "灵活扩展满足业务需求")
    ]
    
    # 创建价值矩阵
    for i, (title, benefit1, benefit2) in enumerate(business_value):
        row = i // 2
        col = i % 2
        
        # 价值卡片背景
        value_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + col * 6.5), Inches(1.5 + row * 2.5), Inches(5.8), Inches(2.2)
        )
        value_bg.fill.solid()
        value_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        value_bg.line.color.rgb = RGBColor(189, 195, 199)
        value_bg.line.width = Pt(2)
        value_bg.shadow.inherit = False
        
        # 价值标题
        value_title = slide.shapes.add_textbox(
            Inches(1 + col * 6.5), Inches(1.7 + row * 2.5), Inches(5.4), Inches(0.6)
        )
        value_title_frame = value_title.text_frame
        value_title_frame.text = title
        value_title_para = value_title_frame.paragraphs[0]
        value_title_para.font.size = Pt(22)
        value_title_para.font.bold = True
        value_title_para.font.color.rgb = success_color
        
        # 价值点1
        value_point1 = slide.shapes.add_textbox(
            Inches(1 + col * 6.5), Inches(2.4 + row * 2.5), Inches(5.4), Inches(0.5)
        )
        value_point1_frame = value_point1.text_frame
        value_point1_frame.text = "• " + benefit1
        value_point1_para = value_point1_frame.paragraphs[0]
        value_point1_para.font.size = Pt(16)
        value_point1_para.font.color.rgb = text_color
        
        # 价值点2
        value_point2 = slide.shapes.add_textbox(
            Inches(1 + col * 6.5), Inches(2.9 + row * 2.5), Inches(5.4), Inches(0.5)
        )
        value_point2_frame = value_point2.text_frame
        value_point2_frame.text = "• " + benefit2
        value_point2_para = value_point2_frame.paragraphs[0]
        value_point2_para.font.size = Pt(16)
        value_point2_para.font.color.rgb = text_color
    
    # 9. 发展规划与路线图
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "发展规划与路线图"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 路线图
    roadmap = [
        ("第一阶段", "核心功能完善", "多模型管理、智能路由、基础监控", "Q1-Q2 2024"),
        ("第二阶段", "企业级增强", "高可用、安全加固、性能优化", "Q3-Q4 2024"),
        ("第三阶段", "生态建设", "插件市场、社区建设、标准化", "Q1-Q2 2025"),
        ("第四阶段", "商业化推广", "SaaS服务、企业定制、行业解决方案", "Q3-Q4 2025")
    ]
    
    # 创建路线图
    for i, (phase, title, desc, timeline) in enumerate(roadmap):
        # 阶段背景
        phase_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 3), Inches(1.8), Inches(2.6), Inches(3.5)
        )
        phase_bg.fill.solid()
        phase_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        phase_bg.line.color.rgb = RGBColor(189, 195, 199)
        phase_bg.line.width = Pt(2)
        phase_bg.shadow.inherit = False
        
        # 阶段名称
        phase_name = slide.shapes.add_textbox(
            Inches(0.9 + i * 3), Inches(2), Inches(2.4), Inches(0.6)
        )
        phase_name_frame = phase_name.text_frame
        phase_name_frame.text = phase
        phase_name_para = phase_name_frame.paragraphs[0]
        phase_name_para.font.size = Pt(18)
        phase_name_para.font.bold = True
        phase_name_para.font.color.rgb = accent_color
        phase_name_para.alignment = PP_ALIGN.CENTER
        
        # 阶段标题
        phase_title = slide.shapes.add_textbox(
            Inches(0.9 + i * 3), Inches(2.7), Inches(2.4), Inches(0.6)
        )
        phase_title_frame = phase_title.text_frame
        phase_title_frame.text = title
        phase_title_para = phase_title_frame.paragraphs[0]
        phase_title_para.font.size = Pt(20)
        phase_title_para.font.bold = True
        phase_title_para.font.color.rgb = primary_color
        phase_title_para.alignment = PP_ALIGN.CENTER
        
        # 阶段描述
        phase_desc = slide.shapes.add_textbox(
            Inches(0.9 + i * 3), Inches(3.4), Inches(2.4), Inches(1.2)
        )
        phase_desc_frame = phase_desc.text_frame
        phase_desc_frame.text = desc
        phase_desc_para = phase_desc_frame.paragraphs[0]
        phase_desc_para.font.size = Pt(16)
        phase_desc_para.font.color.rgb = text_color
        phase_desc_para.alignment = PP_ALIGN.CENTER
        phase_desc_para.word_wrap = True
        
        # 时间线
        timeline_text = slide.shapes.add_textbox(
            Inches(0.9 + i * 3), Inches(4.7), Inches(2.4), Inches(0.5)
        )
        timeline_frame = timeline_text.text_frame
        timeline_frame.text = timeline
        timeline_para = timeline_frame.paragraphs[0]
        timeline_para.font.size = Pt(14)
        timeline_para.font.color.rgb = RGBColor(128, 128, 128)
        timeline_para.alignment = PP_ALIGN.CENTER
    
    # 10. 总结与展望
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.gradient()
    background.fill.gradient.gradient_stops[0].color.rgb = RGBColor(236, 240, 241)
    background.fill.gradient.gradient_stops[1].color.rgb = RGBColor(189, 195, 199)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "总结与展望"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # 总结内容
    summary_text = """
    AICoreDirector作为企业级AI能力管理中枢，通过统一管理、智能调度、灵活扩展、
    成本可控四大核心能力，为企业AI应用提供强有力的支撑。
    
    未来将继续完善功能、优化性能、扩大生态，成为企业AI基础设施的
    重要组成部分。
    """
    
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(2.5))
    summary_frame = summary_box.text_frame
    summary_frame.text = summary_text
    summary_para = summary_frame.paragraphs[0]
    summary_para.font.size = Pt(24)
    summary_para.font.color.rgb = text_color
    summary_para.alignment = PP_ALIGN.CENTER
    summary_para.word_wrap = True
    
    # 联系方式
    contact_text = "感谢聆听 | 欢迎交流与合作"
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = contact_text
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(28)
    contact_para.font.color.rgb = secondary_color
    contact_para.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """主函数"""
    print("正在生成增强版AICoreDirector演示文稿...")
    
    # 创建增强版PPT
    prs = create_enhanced_ppt()
    
    # 保存文件
    output_file = "docs/AICoreDirector_增强版.pptx"
    prs.save(output_file)
    
    print(f"演示文稿已生成: {output_file}")
    print("特点:")
    print("- 丰富的视觉元素和图表")
    print("- 专业的主题配色方案")
    print("- 清晰的信息层次结构")
    print("- 现代化的设计风格")

if __name__ == "__main__":
    main()



