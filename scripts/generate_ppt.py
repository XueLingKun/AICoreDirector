import os
import re
from glob import glob
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


MARKDOWN_PATH = os.path.join("docs", "AICoreDirector_20min_讲稿.md")
OUTPUT_PPTX_PATH = os.path.join("docs", "AICoreDirector_20min_讲稿.pptx")
OUTPUT_PPTX_EXEC_PATH = os.path.join("docs", "AICoreDirector_高管版.pptx")

# Visual Theme
PRIMARY = RGBColor(25, 118, 210)      # Blue 700
ACCENT = RGBColor(2, 136, 209)        # Light Blue 700
TEXT_DARK = RGBColor(33, 33, 33)
TEXT_MUTED = RGBColor(97, 97, 97)
BG_LIGHT = RGBColor(250, 250, 250)
FONT_NAME = "Microsoft YaHei"


def read_markdown(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def split_sections(md_text: str):
    # Sections start with lines beginning with '### '
    lines = md_text.splitlines()
    sections = []
    current_title = None
    current_content = []

    for line in lines:
        if line.startswith("### "):
            if current_title is not None:
                sections.append({"title": current_title, "content": current_content})
            current_title = line[4:].strip()
            current_content = []
        else:
            current_content.append(line)

    if current_title is not None:
        sections.append({"title": current_title, "content": current_content})

    return sections


def add_cover_slide(prs: Presentation, title: str, subtitle: str | None = None):
    layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    # Style title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
    if subtitle is None:
        subtitle = "AI基础能力中枢：统一接入、智能调度、统一管控"
    slide.placeholders[1].text = f"{subtitle}\n{datetime.now().strftime('%Y-%m-%d')}"
    for p in slide.placeholders[1].text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_MUTED

    # Add top-right small placeholder for LOGO
    left = prs.slide_width - Inches(2)
    top = Inches(0.2)
    width = Inches(1.7)
    height = Inches(0.6)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 -> Rectangle
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "LOGO（可替换图片）"
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)

    # Add bottom accent bar
    bar = slide.shapes.add_shape(1, Inches(0), prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_content_slide(prs: Presentation, title: str, content_lines: list[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    def add_bullet(text: str, level: int = 0):
        p = body.add_paragraph() if len(body.paragraphs) > 0 else body.paragraphs[0]
        p.text = text
        p.level = level
        p.font.name = FONT_NAME
        p.font.size = Pt(18 if level == 0 else 16)
        p.font.color.rgb = TEXT_DARK if level == 0 else TEXT_MUTED

    # Parse bullets: lines starting with '- ' (level 0) or '  - ' (level 1) etc.
    for raw in content_lines:
        line = raw.rstrip()
        if not line:
            continue
        m = re.match(r"^(\s*)-\s+(.*)$", line)
        if m:
            indent, text = m.groups()
            level = len(indent) // 2
            add_bullet(text.strip(), level)
        else:
            # Normal paragraph
            add_bullet(line.strip(), 0)

    # Left accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(1.0), Inches(0.15), prs.slide_height - Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_image_slide(prs: Presentation, title: str, image_path: str):
    if not os.path.exists(image_path):
        return
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    # Calculate max area preserving aspect ratio
    max_width = prs.slide_width - Inches(1)
    max_height = prs.slide_height - Inches(2)

    # Let python-pptx scale automatically if we give only width or height
    left = Inches(0.5)
    top = Inches(1.2)
    slide.shapes.add_picture(image_path, left, top, width=max_width)


def add_placeholder_slide(prs: Presentation, title: str, note: str):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
    left = Inches(1.0)
    top = Inches(1.5)
    width = prs.slide_width - Inches(2.0)
    height = prs.slide_height - Inches(3.0)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = note
    p.font.size = Pt(18)
    p.font.name = FONT_NAME
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 248, 248)
    shape.line.color.rgb = RGBColor(180, 180, 180)


def add_bar_chart_slide(prs: Presentation, title: str, categories: list[str], series: list[tuple[str, list[float]]]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)

    x, y, cx, cy = Inches(1), Inches(1.5), prs.slide_width - Inches(2), prs.slide_height - Inches(2.5)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = False


def add_donut_chart_slide(prs: Presentation, title: str, categories: list[str], values: list[float]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("占比", values)

    x, y, cx, cy = Inches(1), Inches(1.5), prs.slide_width - Inches(2), prs.slide_height - Inches(2.5)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.plots[0].has_data_labels = True


def add_tiles_slide(prs: Presentation, title: str, tiles: list[tuple[str, str]]):
    """Create a slide with 2x2 or 1x4 tiles (title + subtitle)."""
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    cols = 2
    rows = (len(tiles) + 1) // 2
    tile_w = (prs.slide_width - Inches(2.5)) / cols
    tile_h = (prs.slide_height - Inches(3.0)) / max(rows, 1)
    left0 = Inches(1.0)
    top0 = Inches(1.5)

    for i, (ttl, sub) in enumerate(tiles):
        r = i // cols
        c = i % cols
        left = left0 + tile_w * c + Inches(0.25)
        top = top0 + tile_h * r + Inches(0.15)
        shape = slide.shapes.add_shape(1, left, top, tile_w - Inches(0.5), tile_h - Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_LIGHT
        shape.line.color.rgb = ACCENT
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = ttl
        p1.font.name = FONT_NAME
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.name = FONT_NAME
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED


def add_value_prop_slide(prs: Presentation, headline: str, metrics: list[tuple[str, str]]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = headline
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    # 3 metrics side-by-side
    count = len(metrics)
    width = (prs.slide_width - Inches(2.5)) / max(count, 1)
    left0 = Inches(1.0)
    top = Inches(2.0)
    for i, (k, v) in enumerate(metrics):
        left = left0 + i * width + Inches(0.25)
        shape = slide.shapes.add_shape(1, left, top, width - Inches(0.5), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_LIGHT
        shape.line.color.rgb = PRIMARY
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = v
        p1.font.name = FONT_NAME
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p2 = tf.add_paragraph()
        p2.text = k
        p2.font.name = FONT_NAME
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def add_before_after_slide(prs: Presentation, title: str, labels: list[str], before: list[float], after: list[float]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Before", before)
    chart_data.add_series("After", after)

    x, y, cx, cy = Inches(1), Inches(1.8), prs.slide_width - Inches(2), prs.slide_height - Inches(3)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False


def add_timeline_slide(prs: Presentation, title: str, milestones: list[tuple[str, str]]):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT_NAME
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    # Horizontal timeline
    y = Inches(3)
    x0 = Inches(1)
    x1 = prs.slide_width - Inches(1)
    line = slide.shapes.add_connector(1, x0, y, x1, y)  # straight connector
    line.line.color.rgb = ACCENT
    step = (x1 - x0) / max(len(milestones) - 1, 1)
    for i, (phase, desc) in enumerate(milestones):
        cx = x0 + i * step
        dot = slide.shapes.add_shape(9, cx - Inches(0.1), y - Inches(0.1), Inches(0.2), Inches(0.2))  # 9 ellipse
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY
        dot.line.fill.background()
        # label
        shape = slide.shapes.add_shape(1, cx - Inches(1.2), y - Inches(1.2), Inches(2.4), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_LIGHT
        shape.line.color.rgb = RGBColor(220, 220, 220)
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = phase
        p1.font.name = FONT_NAME
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_NAME
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED


def build_ppt():
    md_text = read_markdown(MARKDOWN_PATH)
    sections = split_sections(md_text)

    prs = Presentation()

    # Cover from first section title if exists, else default
    cover_title = sections[0]["title"] if sections else "AICoreDirector 20分钟介绍讲稿"
    add_cover_slide(prs, title=cover_title)

    # Add content slides (prefer visuals for key sections)
    BULLETS_PER_SLIDE = 7
    def title_has(sec_title: str, key: str) -> bool:
        return key in sec_title

    for idx, sec in enumerate(sections):
        title = sec["title"].strip()
        content = [line for line in sec["content"] if not line.strip().startswith("---")]
        if idx == 0:
            if content:
                add_content_slide(prs, title, content)
            continue
        # Prefer visuals for certain sections
        if title_has(title, "产品全景"):
            add_tiles_slide(
                prs,
                title="产品全景",
                tiles=[
                    ("统一入口API", "对上REST，对下多源适配"),
                    ("模型池与路由", "标签/等级/健康/QPS/成本"),
                    ("插件化与服务注册", "热插拔与动态路由"),
                    ("提示词模板中心", "集中管理/热加载/复用"),
                    ("监控与成本治理", "健康/QPS/命中/成本聚合"),
                    ("前端控制台", "仪表盘/配置/发现/历史"),
                ],
            )
            continue
        if title_has(title, "六大亮点"):
            add_tiles_slide(
                prs,
                title="六大亮点",
                tiles=[
                    ("多模型统一管理", "元数据标准化/热加载"),
                    ("智能路由", "多因素评分/Fallback"),
                    ("插件生态", "@plugin_api/双入口/并发"),
                    ("服务注册", "注册即入池/动态代理"),
                    ("Prompt中心", "模板沉淀/A-B/热更新"),
                    ("可观可管", "QPS/健康/成本/历史"),
                ],
            )
            continue
        if title_has(title, "架构与关键模块"):
            arch_img = os.path.join("docs", "总体架构图1.png")
            if os.path.exists(arch_img):
                add_image_slide(prs, "总体架构图", arch_img)
            add_tiles_slide(
                prs,
                title="关键模块",
                tiles=[
                    ("API服务层", "入口/参数/路由/异常"),
                    ("业务逻辑层", "会话/选择/批并发"),
                    ("LLM池与路由", "状态/选择/限流/健康"),
                    ("插件加载器", "扫描/注册/热加载"),
                    ("Prompt中心", "模板加载/接口/前端"),
                    ("统计与状态", "命中/成本/历史/偏好"),
                ],
            )
            continue
        if title_has(title, "典型应用场景"):
            add_tiles_slide(
                prs,
                title="典型应用场景",
                tiles=[
                    ("企业AI中枢", "统一接入，降耦合"),
                    ("ISV平台化", "插件化交付"),
                    ("多云治理", "成本/合规路由"),
                    ("内容抽取/分类", "模板+插件沉淀"),
                    ("成本敏感业务", "低成本优先"),
                    ("高可用业务", "回退不中断"),
                ],
            )
            continue
        if title_has(title, "差异化与可持续优势"):
            add_tiles_slide(
                prs,
                title="差异化与优势",
                tiles=[
                    ("中枢定位", "对上治理，对下适配"),
                    ("成本-效果平衡", "可配置/可度量"),
                    ("开发者友好", "插件/并发/Prompt中心"),
                    ("运维友好", "全链路可观测"),
                    ("生态开放", "自研+第三方"),
                    ("可演进", "策略/架构扩展"),
                ],
            )
            continue
        if title_has(title, "部署与商业化"):
            add_tiles_slide(
                prs,
                title="部署与商业化",
                tiles=[
                    ("开发/测试", "本地快速起"),
                    ("生产", "容器化/扩展/监控接入"),
                    ("企业版增值", "权限/审计/SLA"),
                    ("咨询与交付", "选型/接入/落地"),
                ],
            )
            continue
        if title_has(title, "路线图"):
            add_timeline_slide(
                prs,
                title="路线图",
                milestones=[
                    ("短期", "服务发现/插件生态/流式并发"),
                    ("中期", "多模态接入/自适应路由"),
                    ("长期", "生态市场/安全审计/多租户"),
                ],
            )
            continue

        # Detect bullets, paginate for other sections
        bullet_lines = [ln for ln in content if re.match(r"^\s*-\s+", ln)]
        if bullet_lines and len(bullet_lines) > BULLETS_PER_SLIDE:
            # Paginate bullets
            chunks = [bullet_lines[i:i+BULLETS_PER_SLIDE] for i in range(0, len(bullet_lines), BULLETS_PER_SLIDE)]
            for page, chunk in enumerate(chunks, start=1):
                add_content_slide(prs, f"{title}（{page}/{len(chunks)}）", chunk)
        else:
            add_content_slide(prs, title, content)

    # Architecture image (if present)
    arch_img = os.path.join("docs", "总体架构图1.png")
    if os.path.exists(arch_img):
        add_image_slide(prs, "总体架构", arch_img)

    # Mermaid flowcharts (if present at project root)
    for img in sorted(glob(os.path.join(".", "mermaid_*.png"))):
        add_image_slide(prs, f"流程图 {os.path.basename(img)}", img)

    # Charts (illustrative)
    add_donut_chart_slide(prs, "能力来源占比（示意）", ["LLM", "自研AI组件", "传统ML服务"], [50, 30, 20])
    add_bar_chart_slide(prs, "治理指标（示意）", ["成本", "延迟", "错误率", "QPS"], [("当前", [60, 40, 10, 80])])

    # Frontend screenshots placeholder
    add_placeholder_slide(prs, "前端界面概览（占位）", "请在此处替换为前端截图：Dashboard、LLM配置、插件管理、服务发现等页面")

    # Save
    os.makedirs(os.path.dirname(OUTPUT_PPTX_PATH), exist_ok=True)
    prs.save(OUTPUT_PPTX_PATH)
    print(f"PPT 已生成: {OUTPUT_PPTX_PATH}")


def build_exec_ppt():
    prs = Presentation()
    add_cover_slide(prs, title="AICoreDirector 高管版概览")

    # Value prop + metrics
    add_value_prop_slide(
        prs,
        headline="统一AI能力中枢：更快交付、更低成本、更高稳定",
        metrics=[
            ("交付提速", "x3"),
            ("成本下降", "-40%"),
            ("可用性", ">99.9%"),
        ],
    )

    # Pillars as tiles
    add_tiles_slide(
        prs,
        title="平台四大支柱",
        tiles=[
            ("统一入口", "对上统一API，对下多源适配"),
            ("智能路由", "成本/延迟/效果自适应"),
            ("插件/服务生态", "自研AI组件与传统ML统一治理"),
            ("监控与成本治理", "QPS/健康/命中/成本可观可管"),
        ],
    )

    # Capability mix (donut)
    add_donut_chart_slide(prs, "能力来源占比", ["LLM", "自研AI组件", "传统ML服务"], [50, 30, 20])

    # Before vs After (ROI model)
    add_before_after_slide(prs, "投入产出对比（示意）", ["开发成本", "调用成本", "维护成本"], [100, 100, 100], [60, 70, 40])

    # KPI dashboard (bar)
    add_bar_chart_slide(prs, "关键治理指标（示意）", ["成本", "延迟", "错误率", "QPS"], [("当前", [60, 40, 10, 80])])

    # Architecture image or placeholder
    arch_img = os.path.join("docs", "总体架构图1.png")
    if os.path.exists(arch_img):
        add_image_slide(prs, "总体架构图", arch_img)
    else:
        add_placeholder_slide(prs, "总体架构图", "请替换为简化版架构图")

    # Roadmap timeline
    add_timeline_slide(
        prs,
        "路线图",
        [
            ("短期", "服务发现/插件生态，流式与并发治理"),
            ("中期", "多模态接入，自适应路由进化"),
            ("长期", "生态市场化，安全审计，多租户国际化"),
        ],
    )

    # Deployment options
    add_tiles_slide(
        prs,
        title="部署形态",
        tiles=[
            ("本地化/私有化", "容器化部署与运维接入"),
            ("混合云", "多云路由与合规治理"),
            ("平台化输出", "面向合作伙伴与生态对接"),
            ("企业版增值", "权限/审计/SLA/专属支持"),
        ],
    )

    # Case placeholder
    add_placeholder_slide(prs, "典型案例（占位）", "可替换为1-2个行业/业务案例与量化效果")

    os.makedirs(os.path.dirname(OUTPUT_PPTX_EXEC_PATH), exist_ok=True)
    prs.save(OUTPUT_PPTX_EXEC_PATH)
    print(f"PPT 已生成: {OUTPUT_PPTX_EXEC_PATH}")


if __name__ == "__main__":
    build_ppt()
    build_exec_ppt()


